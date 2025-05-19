from transformers import pipeline, AutoTokenizer, AutoModelForCausalLM
import torch
import base64
import io
from PIL import Image
import PyPDF2
import docx
import pptx
import openpyxl
import csv
import json
import zipfile
import os
import mimetypes
from typing import Dict, List, Tuple, Optional
import easyocr
import tempfile

# System prompt - Enhanced for legal assistant
SYSTEM_PROMPT = "You are Bob, a knowledgeable legal AI assistant. Respond conversationally and concisely. Provide helpful legal guidance while always reminding users that this is general information, not legal advice, and they should consult with a licensed attorney for specific situations."

# Pre-load tokenizer and model globally
MODEL_PATH = r"C:\Program Files\Bob-the-lawyer-model\tinyllama_model"
tokenizer = AutoTokenizer.from_pretrained(MODEL_PATH)
model = AutoModelForCausalLM.from_pretrained(MODEL_PATH)

# Check and set device
device = 0 if torch.cuda.is_available() else -1
print(f"Using {'CUDA' if device == 0 else 'CPU'} for inference")

# Initialize the pipeline once
chat_pipeline = pipeline(
    "text-generation",
    model=model,
    tokenizer=tokenizer,
    device=device,
)

# Initialize OCR reader for image text extraction
try:
    ocr_reader = easyocr.Reader(['en'])
    OCR_AVAILABLE = True
    print("OCR capabilities loaded successfully")
except Exception as e:
    print(f"OCR not available: {e}")
    OCR_AVAILABLE = False

class FileProcessor:
    """Handle different file types and extract relevant content"""
    
    @staticmethod
    def extract_text_from_pdf(file_content: bytes) -> str:
        """Extract text from PDF files"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            return f"Error reading PDF: {str(e)}"
    
    @staticmethod
    def extract_text_from_docx(file_content: bytes) -> str:
        """Extract text from Word documents"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = docx.Document(doc_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            return f"Error reading Word document: {str(e)}"
    
    @staticmethod
    def extract_text_from_pptx(file_content: bytes) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            ppt_file = io.BytesIO(file_content)
            prs = pptx.Presentation(ppt_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            return f"Error reading PowerPoint: {str(e)}"
    
    @staticmethod
    def extract_text_from_excel(file_content: bytes) -> str:
        """Extract text from Excel files"""
        try:
            excel_file = io.BytesIO(file_content)
            workbook = openpyxl.load_workbook(excel_file)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    if row_text:
                        text += " | ".join(row_text) + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"
    
    @staticmethod
    def extract_text_from_csv(file_content: bytes) -> str:
        """Extract text from CSV files"""
        try:
            csv_text = file_content.decode('utf-8')
            csv_file = io.StringIO(csv_text)
            reader = csv.reader(csv_file)
            text = ""
            for row in reader:
                text += " | ".join(row) + "\n"
            return text.strip()
        except Exception as e:
            return f"Error reading CSV file: {str(e)}"
    
    @staticmethod
    def extract_text_from_txt(file_content: bytes) -> str:
        """Extract text from plain text files"""
        try:
            encodings = ['utf-8', 'latin-1', 'cp1252']
            for encoding in encodings:
                try:
                    return file_content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            return "Unable to decode text file"
        except Exception as e:
            return f"Error reading text file: {str(e)}"
    
    @staticmethod
    def extract_text_from_image(file_content: bytes) -> str:
        """Extract text from images using OCR"""
        if not OCR_AVAILABLE:
            return "OCR not available. Cannot extract text from image."
        
        try:
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(file_content))
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Save to temporary file for OCR
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp_file:
                image.save(tmp_file.name)
                
                # Perform OCR
                result = ocr_reader.readtext(tmp_file.name)
                
                # Extract text from OCR results
                extracted_text = ""
                for (bbox, text, confidence) in result:
                    if confidence > 0.5:  # Only include text with reasonable confidence
                        extracted_text += text + " "
                
                # Clean up temporary file
                os.unlink(tmp_file.name)
                
                return extracted_text.strip() if extracted_text.strip() else "No text found in image"
                
        except Exception as e:
            return f"Error extracting text from image: {str(e)}"
    
    @staticmethod
    def analyze_file_content(file_content: bytes, file_type: str, filename: str) -> Tuple[str, str]:
        """
        Analyze file content and extract text/information
        
        Returns:
            Tuple[str, str]: (extracted_text, file_analysis)
        """
        file_analysis = f"Analyzing file: {filename} (Type: {file_type})\n"
        extracted_text = ""
        
        try:
            # Document files
            if file_type == 'application/pdf':
                extracted_text = FileProcessor.extract_text_from_pdf(file_content)
                file_analysis += "✓ PDF document processed - legal text may be present\n"
                
            elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                             'application/msword']:
                extracted_text = FileProcessor.extract_text_from_docx(file_content)
                file_analysis += "✓ Word document processed - may contain contracts or legal documents\n"
                
            elif file_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                             'application/vnd.ms-powerpoint']:
                extracted_text = FileProcessor.extract_text_from_pptx(file_content)
                file_analysis += "✓ PowerPoint presentation processed\n"
                
            elif file_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                             'application/vnd.ms-excel']:
                extracted_text = FileProcessor.extract_text_from_excel(file_content)
                file_analysis += "✓ Excel spreadsheet processed\n"
                
            elif file_type == 'text/csv':
                extracted_text = FileProcessor.extract_text_from_csv(file_content)
                file_analysis += "✓ CSV file processed\n"
                
            elif file_type == 'text/plain':
                extracted_text = FileProcessor.extract_text_from_txt(file_content)
                file_analysis += "✓ Text file processed\n"
                
            # Image files
            elif file_type.startswith('image/'):
                extracted_text = FileProcessor.extract_text_from_image(file_content)
                file_analysis += "✓ Image processed with OCR - may contain scanned legal documents\n"
                
            # JSON files
            elif file_type == 'application/json':
                try:
                    json_data = json.loads(file_content.decode('utf-8'))
                    extracted_text = json.dumps(json_data, indent=2)
                    file_analysis += "✓ JSON file processed\n"
                except:
                    extracted_text = "Error parsing JSON file"
                    
            # Archive files
            elif file_type in ['application/zip', 'application/x-zip-compressed']:
                try:
                    zip_file = zipfile.ZipFile(io.BytesIO(file_content))
                    file_list = zip_file.namelist()
                    extracted_text = f"Archive contains {len(file_list)} files:\n" + "\n".join(file_list)
                    file_analysis += f"✓ ZIP archive processed - contains {len(file_list)} files\n"
                except:
                    extracted_text = "Error reading ZIP archive"
                    
            else:
                extracted_text = f"Unsupported file type: {file_type}"
                file_analysis += f"⚠ Unsupported file type: {file_type}\n"
            
            # Analyze extracted text for legal keywords
            if extracted_text and not extracted_text.startswith("Error"):
                legal_keywords = [
                    'contract', 'agreement', 'terms', 'conditions', 'clause', 'liability',
                    'indemnification', 'warranty', 'confidentiality', 'non-disclosure',
                    'intellectual property', 'copyright', 'trademark', 'patent',
                    'jurisdiction', 'governing law', 'dispute resolution', 'arbitration',
                    'termination', 'breach', 'damages', 'penalty', 'force majeure'
                ]
                
                found_keywords = []
                text_lower = extracted_text.lower()
                for keyword in legal_keywords:
                    if keyword in text_lower:
                        found_keywords.append(keyword)
                
                if found_keywords:
                    file_analysis += f"🏛 Legal keywords detected: {', '.join(found_keywords[:5])}"
                    if len(found_keywords) > 5:
                        file_analysis += f" and {len(found_keywords) - 5} more"
                    file_analysis += "\n"
                
        except Exception as e:
            extracted_text = f"Error processing file: {str(e)}"
            file_analysis += f"❌ Error processing file: {str(e)}\n"
        
        return extracted_text, file_analysis

def generate_reply(user_input: str,
                    max_new_tokens: int = 150,
                    temperature: float = 0.7,
                    top_p: float = 0.9,
                    file_attachments: Optional[List[Dict]] = None) -> str:
    """
    Generates a chat-style reply using the pre-built text-generation pipeline.
    Now supports comprehensive file attachments including documents, images, and camera captures.

    Args:
        user_input (str): The user's message
        max_new_tokens (int): Max tokens to generate
        temperature (float): Sampling temperature
        top_p (float): Nucleus sampling probability
        file_attachments (List[Dict]): List of file attachments with format:
            [{'filename': str, 'content': bytes, 'mime_type': str}, ...]

    Returns:
        str: The assistant's reply with file analysis
    """
    
    # Process file attachments if present
    file_context = ""
    extracted_content = ""
    
    if file_attachments:
        file_context = f"\n\n📎 I can see you've attached {len(file_attachments)} file(s):\n"
        
        for i, attachment in enumerate(file_attachments, 1):
            filename = attachment.get('filename', f'file_{i}')
            content = attachment.get('content', b'')
            mime_type = attachment.get('mime_type', 'application/octet-stream')
            
            # Guess mime type if not provided
            if mime_type == 'application/octet-stream':
                mime_type, _ = mimetypes.guess_type(filename)
                mime_type = mime_type or 'application/octet-stream'
            
            # Process the file
            extracted_text, file_analysis = FileProcessor.analyze_file_content(content, mime_type, filename)
            
            file_context += f"\n{i}. {filename}\n{file_analysis}"
            
            # Add extracted content to context (truncate if too long)
            if extracted_text and not extracted_text.startswith("Error"):
                content_preview = extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text
                extracted_content += f"\n\n--- Content from {filename} ---\n{content_preview}\n"
    
    # Construct the enhanced prompt
    enhanced_input = user_input
    if file_context:
        enhanced_input += file_context
    if extracted_content:
        enhanced_input += "\n\nExtracted content from files:" + extracted_content
        enhanced_input += "\n\nPlease review this content and provide legal guidance."
    
    # Construct the prompt with legal context
    prompt = (
        f"<|system|> {SYSTEM_PROMPT}\n"
        f"<|user|> {enhanced_input}\n"
        f"<|assistant|>"
    )
    
    # Generate response
    try:
        outputs = chat_pipeline(
            prompt,
            max_new_tokens=max_new_tokens,
            do_sample=True,
            temperature=temperature,
            top_p=top_p,
            pad_token_id=tokenizer.eos_token_id,
            eos_token_id=tokenizer.eos_token_id,
        )
        
        # Extract the assistant's response
        generated_text = outputs[0]['generated_text']
        reply = generated_text.split("<|assistant|>")[-1].strip()
        
        # Add file processing summary if files were attached
        if file_attachments:
            reply = f"📋 File Analysis Complete:\n{file_context}\n\n{reply}"
        
        # Add legal disclaimer
        if not "disclaimer" in reply.lower():
            reply += "\n\n⚖ *Legal Disclaimer*: This analysis provides general information only, not legal advice. Please consult with a licensed attorney for your specific situation, especially for document review and legal interpretation."
        
        return reply
        
    except Exception as e:
        error_msg = f"I apologize, but I encountered an error processing your request: {str(e)}. Please try again or contact support."
        return error_msg

# Camera capture function (for integration with GUI applications)
def process_camera_capture(image_data: bytes, question: str = "Please analyze this image for legal content") -> str:
    """
    Process camera capture and analyze for legal content
    
    Args:
        image_data (bytes): Raw image data from camera
        question (str): User's question about the image
    
    Returns:
        str: Analysis result
    """
    attachment = {
        'filename': 'camera_capture.jpg',
        'content': image_data,
        'mime_type': 'image/jpeg'
    }
    
    return generate_reply(question, file_attachments=[attachment])

# Example usage function
def example_usage():
    """Example of how to use the enhanced file processing capabilities"""
    
    # Example 1: Text input only
    response1 = generate_reply("What should I look for in an employment contract?")
    print("Response 1:", response1)
    
    # Example 2: With file attachment (simulated)
    # In real usage, you would read the file content as bytes
    # file_content = open('contract.pdf', 'rb').read()
    # attachments = [{'filename': 'contract.pdf', 'content': file_content, 'mime_type': 'application/pdf'}]
    # response2 = generate_reply("Please review this contract", file_attachments=attachments)
    # print("Response 2:", response2)

if __name__ == "__main__":
    print("Enhanced Legal AI Assistant with File Support Loaded!")
    print("Supported file types:")
    print("- Documents: PDF, Word (.docx), PowerPoint (.pptx), Excel (.xlsx)")
    print("- Text files: TXT, CSV, JSON")
    print("- Images: JPG, PNG, GIF, BMP, WEBP (with OCR)")
    print("- Archives: ZIP files")
    print("- Camera captures: Supported through process_camera_capture()")
    
    # Run example if needed
    # example_usage()
