import flet as ft 
import sqlite3    
from datetime import datetime 
from model_handler import generate_reply  
from sidebar import render_sidebar  
from about_us import create_about_us_view
import os           
import platform    
import requests     

import pypdf  
from docx import Document  
from pptx import Presentation 
import openpyxl  


class LawyerChatBotApp:

    def __init__(self, page: ft.Page):
        self.page = page  
        self.page.theme_mode = ft.ThemeMode.LIGHT 
        self.chat = ft.ListView(expand=True, spacing=10, auto_scroll=True)  
        self.current_discussion = None
        self.main_content_area = ft.Container(expand=True)
        self.chat_view = None

        self.theme_toggle = ft.IconButton(
            icon=ft.Icons.DARK_MODE,  
            on_click=self.toggle_theme,  
            tooltip="Toggle dark/light mode",  
        )
        
        self.file_picker = ft.FilePicker(
            on_result=self.handle_file_upload,  
        )
        self.page.overlay.append(self.file_picker)  
        self.sidebar = render_sidebar(self) if 'render_sidebar' in globals() else ft.Container() 
        
        self.user_input = ft.TextField(
            hint_text="Type your legal question or upload documents...",
            expand=True,  
            multiline=True, 
            min_lines=1,    
            max_lines=5,     
        )
        self.send_button = ft.IconButton(
            icon=ft.Icons.SEND,  
            on_click=self.send_click, 
            tooltip="Send question", 
        )
        self.upload_button = ft.IconButton(
            icon=ft.Icons.UPLOAD_FILE,  
            on_click=self.upload_files,   
            tooltip="Upload documents",   
        )
        self.search_button = ft.IconButton(
            icon=ft.Icons.SEARCH,  
            on_click=self.web_search_click,  
            tooltip="Search the web",        
        )
        self.current_files = []  
        self.initialize_database()  
        self.init_ui()              
        self.update_theme_colors()  
        self.switch_discussion(self.current_discussion) 


    def init_ui(self):

        self.page.title = "Bob the lawyer"  
        self.page.window_width = 900        
        self.page.window_height = 700       
        self.page.padding = 20              

        header = ft.Row(
            controls=[
                ft.Text("Chat", size=24, weight=ft.FontWeight.BOLD), 
                self.theme_toggle,  
            ],
            alignment=ft.MainAxisAlignment.SPACE_BETWEEN,
        )

        self.chat_view = ft.Column(
            expand=True,  
            controls=[
                header,  
                ft.Divider(color=ft.Colors.GREY_600 if self.page.theme_mode == ft.ThemeMode.DARK else ft.Colors.GREY_300),
                self.chat,  
                ft.Row(
                    [
                        self.upload_button, 
                        self.user_input,     
                        self.search_button,  
                        self.send_button,    
                    ],
                    spacing=5,  
                    alignment=ft.MainAxisAlignment.START,  
                ),
            ],
        )

        self.main_content_area.content = self.chat_view

        content = ft.Row(
            controls=[
                self.sidebar,  
                ft.Container(width=1, bgcolor=ft.Colors.GREY_600 if self.page.theme_mode == ft.ThemeMode.DARK else ft.Colors.GREY_300), 
                self.main_content_area,
            ],
            expand=True,  
        )

        self.page.add(content)  
        self.user_input.focus() 


    def initialize_database(self):
        """Connect to the database and set up the cursor."""
        self.conn = sqlite3.connect(self.get_database_path())
        self.cursor = self.conn.cursor()


    def _create_table_if_not_exists(self, table_name):

        if not table_name:
            print("Error: Attempted to create a table with no name.")
            return
        try:
            self.cursor.execute(f'''
            CREATE TABLE IF NOT EXISTS "{table_name}" (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                sender TEXT NOT NULL,
                message TEXT NOT NULL,
                timestamp DATETIME NOT NULL
            )
        ''')
            
        except sqlite3.Error as e:
            print(f"Error creating table {table_name}: {e}")


    def get_database_path(self):
        system = platform.system()  
        if system == "Windows":  
            app_data_dir = os.path.join(os.environ['LOCALAPPDATA'], "BobTheLawyer")  
        elif system == "Darwin":  
            app_data_dir = os.path.expanduser("~/Library/Application Support/BobTheLawyer")  
        else:  
            app_data_dir = os.path.expanduser("~/.bobthelawyer")  

        os.makedirs(app_data_dir, exist_ok=True)  
        db_path = os.path.join(app_data_dir, "database.db")  
        return db_path


    def load_previous_messages(self, table_to_load: str):
        self.clear_chat()  
        if table_to_load:  

            try:
                self.cursor.execute(
                    f'SELECT sender, message FROM "{table_to_load}" ORDER BY timestamp' 
                )
                messages = self.cursor.fetchall() 

                for sender, message_content in messages: 
                    if sender == "user":  
                        self.chat.controls.append(self.create_user_message(message_content))  
                    elif sender == "file":  
                        parts = message_content.split(": ", 1)  
                        if len(parts) == 2:  
                            file_name, content_preview = parts  
                            self.chat.controls.append(self.create_file_message(file_name, content_preview))  
                        else:
                            self.chat.controls.append(self.create_bot_message(f"File (error displaying): {message_content}")) 
                    elif sender == "bot" or sender == "system": 
                        self.chat.controls.append(self.create_bot_message(message_content)) 
                    else: 
                        self.chat.controls.append(self.create_bot_message(f"Unknown ({sender}): {message_content}")) 
                self.page.update()  

            except sqlite3.OperationalError as e:  
                print(f"Error loading messages from table '{table_to_load}': {str(e)}")
                self.chat.controls.append(
                    self.create_bot_message(f"⚠️ Error loading discussion '{table_to_load}': {str(e)}")
                )
                self.page.update() 


    def clear_chat(self):
        self.chat.controls.clear()
        self.current_files = []
        self.page.update()


    def switch_discussion(self, table_name):
        self.current_discussion = table_name  
        if table_name is None:  
            self.clear_chat() 
            self.chat.controls.append( 
                ft.Container(
                    content=ft.Text(
                        "Select or create a discussion from the sidebar to begin.",
                        italic=True,
                        opacity=0.6,
                        text_align=ft.TextAlign.CENTER
                    ),
                    alignment=ft.alignment.center,
                    padding=20
                )
            )
            self.user_input.disabled = True 
            self.send_button.disabled = True 
            self.upload_button.disabled = True 
            self.search_button.disabled = True 
        else: 
            self.load_previous_messages(table_name)  
            self.user_input.disabled = False 
            self.send_button.disabled = False 
            self.upload_button.disabled = False 
            self.search_button.disabled = False 
        self.page.update() 


    def show_about_us_view(self):
        photos_dir = "./photos"  # Relative path
        simulated_member_images = []
        for filename in ["McBright.jpeg", "Armel.jpeg", "Royce.jpeg", "Rejoice.jpeg", "MIKe.jpg","Group10.jpeg"]:
            image_path = os.path.join(photos_dir, filename)  # Construct path
            if os.path.exists(image_path):  # Check if file exists
                simulated_member_images.append(image_path)
            else:
                simulated_member_images.append(filename)  # Fallback to assets

        about_us_content = create_about_us_view(self, simulated_member_images) 
        self.main_content_area.content = about_us_content
        
        self.page.update()


    def go_back_to_main(self):
        self.main_content_area.content = self.chat_view
        self.page.update()


    def create_user_message(self, message):
        
        is_dark = self.page.theme_mode == ft.ThemeMode.DARK  
        text_width = min(600, len(message) * 8 + 100) 
        return ft.Row(
            controls=[
                ft.Container(expand=True),  
                ft.Container(
                    content=ft.Text(
                        value=message,
                        selectable=True,
                        size=15,
                        color=ft.Colors.WHITE if is_dark else ft.Colors.BLACK,
                    ),
                    alignment=ft.alignment.center_right, 
                    bgcolor=ft.Colors.BLUE_800 if is_dark else ft.Colors.BLUE_100, 
                    padding=ft.padding.symmetric(horizontal=14, vertical=10), 
                    border_radius=ft.border_radius.only(
                        top_left=16,
                        top_right=16,
                        bottom_left=16,
                        bottom_right=4,
                    ),
                    margin=ft.margin.only(bottom=6),
                    shadow=ft.BoxShadow( 
                        spread_radius=0.5,
                        blur_radius=3,
                        color=ft.Colors.with_opacity(0.1, ft.Colors.BLACK),
                        offset=ft.Offset(0, 1.5),
                    ),
                    width=text_width, 
                )
            ],
            alignment=ft.MainAxisAlignment.END,  
            expand=True, 
        )


    def create_bot_message(self, message):
        is_dark = self.page.theme_mode == ft.ThemeMode.DARK 
        text_width = min(600, len(message) * 8 + 100) 
        return ft.Row(
            controls=[
                ft.Container(
                    content=ft.Column(
                        controls=[
                            ft.Text(
                                "BOB:",
                                size=15,
                                weight=ft.FontWeight.BOLD,
                                color=ft.Colors.BLUE_200 if is_dark else ft.Colors.BLUE_800,
                            ),
                            ft.Container(
                                content=ft.Text(
                                    message,
                                    selectable=True,
                                    size=15,
                                    color=ft.Colors.WHITE if is_dark else ft.Colors.BLACK,
                                ),
                            )
                        ],
                        spacing=4, 
                        tight=True,
                        scroll=ft.ScrollMode.AUTO,
                    ),
                    alignment=ft.alignment.center_left, 
                    bgcolor=ft.Colors.GREEN_800 if is_dark else ft.Colors.GREEN_100, 
                    padding=ft.padding.symmetric(horizontal=14, vertical=10), 
                    border_radius=ft.border_radius.only( 
                        top_left=16,
                        top_right=16,
                        bottom_left=4,
                        bottom_right=16,
                    ),
                    margin=ft.margin.only(bottom=6),
                    shadow=ft.BoxShadow(
                        spread_radius=0.5, 
                        blur_radius=3,
                        color=ft.Colors.with_opacity(0.1, ft.Colors.BLACK),
                        offset=ft.Offset(0, 1.5),
                    ),
                    width=text_width,
                )
            ],
            alignment=ft.MainAxisAlignment.START, 
            vertical_alignment=ft.CrossAxisAlignment.START, 
            expand=True, 
            scroll=ft.ScrollMode.AUTO,
        )


    def create_file_message(self, file_name, content_preview):
        is_dark = self.page.theme_mode == ft.ThemeMode.DARK 
        text_width = min(600, (len(file_name) + len(content_preview)) * 6 + 100) 
        return ft.Row(
            [
                ft.Container(
                    content=ft.Column(
                        [
                            ft.Text(
                                f"📄 {file_name}",
                                size=15,
                                weight=ft.FontWeight.BOLD,
                                color=ft.Colors.WHITE if is_dark else ft.Colors.BLUE_800, 
                            ),
                            ft.Text(
                                content_preview,
                                selectable=True,
                                size=12,
                                color=ft.Colors.WHITE if is_dark else ft.Colors.BLACK, 
                            )
                        ],
                        spacing=4, 
                        tight=True, 
                    ),
                    alignment=ft.alignment.center_right, 
                    bgcolor=ft.Colors.BLUE_800 if is_dark else ft.Colors.BLUE_100, 
                    padding=ft.padding.symmetric(horizontal=14, vertical=10), 
                    border_radius=ft.border_radius.only( 
                        top_left=16,
                        top_right=16,
                        bottom_left=16,
                        bottom_right=4,
                    ),
                    margin=ft.margin.only(bottom=6),
                    shadow=ft.BoxShadow( 
                        spread_radius=0.5,
                        blur_radius=3,
                        color=ft.Colors.with_opacity(0.1, ft.Colors.BLACK),
                        offset=ft.Offset(0, 1.5),
                    ),
                    width=text_width, 
                )
            ],
            alignment=ft.MainAxisAlignment.END, 
        )


    def store_message(self, sender, message):
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

        try:
            self.cursor.execute(f'''
                INSERT INTO "{self.current_discussion}" (sender, message, timestamp)
                VALUES (?, ?, ?)
            ''', (sender, message, timestamp)) 
            self.conn.commit() 
        
        except sqlite3.OperationalError as e: 
            print(f"Error storing message (table '{self.current_discussion}' might not exist): {str(e)}")
            
            self._create_table_if_not_exists(self.current_discussion)
            
            try:
                self.cursor.execute(f'''
                    INSERT INTO "{self.current_discussion}" (sender, message, timestamp)
                    VALUES (?, ?, ?)
                ''', (sender, message, timestamp))
                self.conn.commit()
            
            except sqlite3.Error as e_retry:
                print(f"Retry failed for storing message: {e_retry}")


    def upload_files(self, e):
        
        self.file_picker.pick_files(
            allowed_extensions=[
                "pdf", "docx", "pptx", "xls", "xlsx", 
            ],
        )


    def handle_file_upload(self, e: ft.FilePickerResultEvent):
        if not e.files: 
            return
            
        for uploaded_file in e.files: 
            file_path = uploaded_file.path  
            file_name = uploaded_file.name  
            file_ext = os.path.splitext(file_name)[1].lower() 
            
            try:
                if file_ext == '.pdf': 
                    text = self.extract_text_from_pdf(file_path)
                elif file_ext == '.docx':
                    text = self.extract_text_from_docx(file_path)
                elif file_ext == '.pptx':
                    text = self.extract_text_from_pptx(file_path)
                elif file_ext in ('.xls', '.xlsx'):
                    text = self.extract_text_from_excel(file_path)
                else:
                    text = f"Unsupported file type: {file_ext}"
                
                self.current_files.append({
                    "name": file_name,
                    "path": file_path,
                    "content": text
                })
                
                original_lines = text.splitlines() 
                num_original_lines = len(original_lines) 

                if num_original_lines > 4: 
                    preview_lines_to_show = original_lines[:3] 
                    preview = "\n".join(preview_lines_to_show) + "\n..."
                else:
                    preview = "\n".join(original_lines)

                self.chat.controls.append(
                    self.create_file_message(file_name, preview))
                
                self.store_message("user", f"Uploaded document: {file_name}") 
                self.store_message("file", f"{file_name}: {preview}") 
                
            except Exception as ex: 
                error_msg = f"Error processing {file_name}: {str(ex)}"
                self.chat.controls.append(
                    self.create_bot_message(error_msg)) 
                    
        self.page.update() 


    def extract_text_from_pdf(self, file_path: str) -> str:
        text = []
        with open(file_path, 'rb') as file: 
            reader = pypdf.PdfReader(file) 
            for page in reader.pages: 
                text.append(page.extract_text()) 
        return "\n".join(text)


    def extract_text_from_docx(self, file_path: str) -> str:
        doc = Document(file_path) 
        return "\n".join([para.text for para in doc.paragraphs])


    def extract_text_from_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path) 
        text = []
        for slide in prs.slides: 
            for shape in slide.shapes: 
                if hasattr(shape, "text"): 
                    text.append(shape.text)
        return "\n".join(text)


    def extract_text_from_excel(self, file_path: str) -> str:
        
        workbook = openpyxl.load_workbook(file_path, data_only=True) 
        all_text = []
        for sheet_name in workbook.sheetnames: 
            sheet = workbook[sheet_name] 
            sheet_text = []
            for row in sheet.iter_rows():
                for cell in row: 
                    if cell.value is not None: 
                        sheet_text.append(str(cell.value)) 
            if sheet_text: 
                all_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(sheet_text)) 
        return "\n\n".join(all_text)


    def toggle_theme(self, e):
        self.page.theme_mode = (
            ft.ThemeMode.DARK 
            if self.page.theme_mode == ft.ThemeMode.LIGHT
            else ft.ThemeMode.LIGHT 
        )
        self.theme_toggle.icon = (
            ft.Icons.LIGHT_MODE 
            if self.page.theme_mode == ft.ThemeMode.DARK
            else ft.Icons.DARK_MODE 
        )
        self.update_theme_colors() 
        self.sidebar.refresh_sidebar(self.page)

        if self.main_content_area.content is not self.chat_view:
            self.show_about_us_view()
        elif self.current_discussion:
            self.load_previous_messages(self.current_discussion) 
        self.page.update()


    def update_theme_colors(self):
        is_dark = self.page.theme_mode == ft.ThemeMode.DARK
        
        self.page.bgcolor = ft.Colors.GREY_900 if is_dark else ft.Colors.GREY_50
        
        self.user_input.color = ft.Colors.WHITE if is_dark else ft.Colors.BLACK
        self.user_input.bgcolor = ft.Colors.with_opacity(0.1, ft.Colors.GREY_500 if is_dark else ft.Colors.GREY_200)
        self.user_input.border_color = ft.Colors.GREY_600 if is_dark else ft.Colors.GREY_300
        
        self.send_button.Icon_color = ft.Colors.WHITE if is_dark else ft.Colors.BLUE_700
        self.upload_button.Icon_color = ft.Colors.WHITE if is_dark else ft.Colors.BLUE_700
        self.theme_toggle.Icon_color = ft.Colors.WHITE if is_dark else ft.Colors.BLUE_700
        self.search_button.Icon_color = ft.Colors.WHITE if is_dark else ft.Colors.BLUE_700   


    def web_search_click(self, e):
        query = self.user_input.value.strip() 
        if not query:
            return

        if not self.current_discussion:
            print("Please start a new discussion first.")
            return

        self.store_message("user", f"WEB SEARCH: {query}") 
        self.chat.controls.append(self.create_user_message(f"🔍 Searching: {query}")) 

        thinking = ft.Container(
            ft.Row([
                ft.ProgressRing(width=20, height=20, stroke_width=2),
                ft.Text("Searching the web...")
            ], spacing=10),
            alignment=ft.alignment.center_left,
        )
        self.chat.controls.append(thinking)
        self.theme_toggle.disabled = True
        self.user_input.disabled = True
        self.send_button.disabled = True
        self.upload_button.disabled = True
        self.search_button.disabled = True
        self.page.update()

        try:
            api_key = "c993fe8beec5d447b42da49cec429aab6460e9170118ba6eb56473f574705105" 
            search_query = f"{query} from the Law of Cameroon"
            params = {
                "q": search_query,
                "engine": "google",
                "api_key": api_key,
                "num": 2
            }
            
            response = requests.get("https://serpapi.com/search", params=params, timeout=10) 
            response.raise_for_status()  
            results = response.json() 
            
            if "organic_results" in results: 
                display_output = ["🌐 Web Results:"]
                ai_context = "\n\nFROM THE INTERNET I FOUND:\n"
                
                for idx, res in enumerate(results["organic_results"][:2], 1): 
                    title = res.get('title', 'No title')
                    snippet = res.get('snippet', 'No description')
                    link = res.get('link', 'No URL')
                    
                    display_output.append(f"{idx}. {title}\n   {snippet}\n   {link}\n")
                    
                    ai_context += f"\nResult {idx}:\nTitle: {title}\nContent: {snippet}\nSource: {link}\n"
                
                display_result = "\n".join(display_output)
                
                enhanced_query = f"{query}{ai_context}\n\nPlease provide a comprehensive answer based on the search results above."
                
            else:
                display_result = "🔍 No results found"
                enhanced_query = f"{query}\n\nFROM THE INTERNET I FOUND:\nNo relevant search results were found for this query."
            
            self.store_message("bot", display_result) 
            self.chat.controls.remove(thinking) 
            self.chat.controls.append(self.create_bot_message(display_result))
            
            thinking_ai = ft.Container(
                ft.Row([
                    ft.ProgressRing(width=20, height=20, stroke_width=2),
                    ft.Text("Analyzing search results...")
                ], spacing=10),
                alignment=ft.alignment.center_left,
            )
            self.chat.controls.append(thinking_ai)
            self.page.update()
            
            try:
                ai_reply = generate_reply(enhanced_query)
                self.store_message("bot", ai_reply)
                self.chat.controls.remove(thinking_ai)
                self.chat.controls.append(self.create_bot_message(ai_reply))
            except Exception as ai_err:
                ai_error_msg = f"⚠️ Error analyzing results: {str(ai_err)}"
                self.store_message("system", ai_error_msg)
                self.chat.controls.remove(thinking_ai)
                self.chat.controls.append(self.create_bot_message(ai_error_msg))
            
        except Exception as err: 
            error_msg = f"⚠️ Search failed: {str(err)}"
            self.store_message("system", error_msg) 
            self.chat.controls.remove(thinking) 
            self.chat.controls.append(self.create_bot_message(error_msg)) 
            
        self.user_input.value = ""
        self.user_input.disabled = False
        self.send_button.disabled = False
        self.upload_button.disabled = False
        self.search_button.disabled = False
        self.theme_toggle.disabled = False  
        self.page.update()


    def send_click(self, e):
        question = self.user_input.value.strip() 
        if not question and not self.current_files: 
            return
        
        if not self.current_discussion:  
            print("Please start a new discussion first.")
            return

        context = ""
        if self.current_files: 
            context = "\n\n[Attached Files Context]\n"
            for file in self.current_files: 
                context += f"\nFile: {file['name']}\nContent:\n{file['content'][:1000]}\n"
            self.current_files = []  

        if question: 
            full_question = question + context
            self.store_message("user", question) 
            self.chat.controls.append(self.create_user_message(question)) 
        else:  
            full_question = "Please analyze these documents:" + context
            self.store_message("user", "Uploaded documents for analysis") 

        if context and not question: 
            self.chat.controls.append(self.create_bot_message("Received documents for analysis"))

        thinking = ft.Container(
            ft.Row([
                ft.ProgressRing(width=20, height=20, stroke_width=2),
                ft.Text("Thinking...")
            ], spacing=10),
            alignment=ft.alignment.center_left,
        )
        self.chat.controls.append(thinking)

        self.user_input.disabled = True
        self.send_button.disabled = True
        self.upload_button.disabled = True
        self.search_button.disabled = True
        self.theme_toggle.disabled = True 
        self.page.update()

        try:
            reply = generate_reply(full_question) 
            self.store_message("bot", reply) 
        except Exception as err: 
            reply = f"⚠️ Error: {str(err)}"
            self.store_message("system", f"Error: {str(err)}") 

        self.chat.controls.remove(thinking) 
        self.chat.controls.append(self.create_bot_message(reply)) 

        self.user_input.value = ""
        self.user_input.disabled = False
        self.send_button.disabled = False
        self.upload_button.disabled = False
        self.theme_toggle.disabled = False
        self.search_button.disabled = False
        self.page.update() 
        self.user_input.focus() 


    def __del__(self):
        if hasattr(self, 'conn'): 
            self.conn.close() 


def main(page: ft.Page):
    """Main function to start the Flet application."""
    LawyerChatBotApp(page)  

ft.app(target=main, assets_dir="src/assets")
