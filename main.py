import flet as ft
import sqlite3
from datetime import datetime
from model_handler import generate_reply
from sidebar import render_sidebar  # Optional sidebar module if you use one
import tempfile
import os
from typing import Optional
import io

# Import file processing libraries
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract

class LawyerChatBotApp:
    def __init__(self, page: ft.Page):
        self.page = page
        self.page.theme_mode = ft.ThemeMode.LIGHT  # Default to light theme
        self.chat = ft.ListView(expand=True, spacing=10, auto_scroll=True)
        
        # Theme toggle button
        self.theme_toggle = ft.IconButton(
            icon=ft.icons.DARK_MODE,
            on_click=self.toggle_theme,
            tooltip="Toggle dark/light mode",
        )
        
        # Initialize file picker
        self.file_picker = ft.FilePicker(
            on_result=self.handle_file_upload,
        )
        self.page.overlay.append(self.file_picker)
        self.sidebar = render_sidebar() if 'render_sidebar' in globals() else ft.Container()
        
        # Input controls
        self.user_input = ft.TextField(
            hint_text="Type your legal question or upload documents...",
            expand=True,
            multiline=True,
            min_lines=1,
            max_lines=5,
        )
        self.send_button = ft.IconButton(
            icon=ft.icons.SEND,
            on_click=self.send_click,
            tooltip="Send question",
        )
        self.upload_button = ft.IconButton(
            icon=ft.icons.UPLOAD_FILE,
            on_click=self.upload_files,
            tooltip="Upload documents",
        )
        self.current_files = []  # Stores information about uploaded files
        
        # Initialize database and load previous messages
        self.initialize_database()
        self.load_previous_messages()
        self.init_ui()
        self.update_theme_colors()  # Set initial theme colors

    def initialize_database(self):
        """Create database and table if they don't exist"""
        self.conn = sqlite3.connect('database.db')
        self.cursor = self.conn.cursor()
        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS messages (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                sender TEXT NOT NULL,
                message TEXT NOT NULL,
                timestamp DATETIME NOT NULL
            )
        ''')
        self.conn.commit()

    def load_previous_messages(self):
        """Load previous messages from database when app starts"""
        self.cursor.execute('SELECT sender, message FROM messages ORDER BY timestamp')
        messages = self.cursor.fetchall()
        
        for sender, message in messages:
            if sender == "user":
                self.chat.controls.append(self.create_user_message(message))
            else:  # bot or system messages
                self.chat.controls.append(self.create_bot_message(message))
        self.page.update()

    def create_user_message(self, message):
        is_dark = self.page.theme_mode == ft.ThemeMode.DARK
        return ft.Row(
            [
                ft.Container(
                    content=ft.Text(
                        message,
                        selectable=True,
                        size=15,
                        color=ft.colors.WHITE if is_dark else ft.colors.BLACK,
                    ),
                    alignment=ft.alignment.center_right,
                    bgcolor=ft.colors.BLUE_800 if is_dark else ft.colors.BLUE_100,
                    padding=ft.padding.symmetric(horizontal=14, vertical=10),
                    border_radius=ft.border_radius.only(
                        top_left=16,
                        top_right=16,
                        bottom_left=16,
                        bottom_right=4,
                    ),
                    margin=ft.margin.only(bottom=6),
                    shadow=ft.BoxShadow(
                        spread_radius=0.5,
                        blur_radius=3,
                        color=ft.colors.with_opacity(0.1, ft.colors.BLACK),
                        offset=ft.Offset(0, 1.5),
                    ),
                )
            ],
            alignment=ft.MainAxisAlignment.END,
        )

    def create_bot_message(self, message):
        is_dark = self.page.theme_mode == ft.ThemeMode.DARK
        return ft.Row(
            [
                ft.Container(
                    content=ft.Column(
                        [
                            ft.Text(
                                "BOB:",
                                size=15,
                                weight=ft.FontWeight.BOLD,
                                color=ft.colors.BLUE_200 if is_dark else ft.colors.BLUE_800,
                            ),
                            ft.Text(
                                message,
                                selectable=True,
                                size=15,
                                color=ft.colors.WHITE if is_dark else ft.colors.BLACK,
                            )
                        ],
                        spacing=4,
                        tight=True,
                    ),
                    alignment=ft.alignment.center_left,
                    bgcolor=ft.colors.GREEN_800 if is_dark else ft.colors.GREEN_100,
                    padding=ft.padding.symmetric(horizontal=14, vertical=10),
                    border_radius=ft.border_radius.only(
                        top_left=16,
                        top_right=16,
                        bottom_left=4,
                        bottom_right=16,
                    ),
                    margin=ft.margin.only(bottom=6),
                    shadow=ft.BoxShadow(
                        spread_radius=0.5,
                        blur_radius=3,
                        color=ft.colors.with_opacity(0.1, ft.colors.BLACK),
                        offset=ft.Offset(0, 1.5),
                    ),
                )
            ],
            alignment=ft.MainAxisAlignment.START,
        )

    def create_file_message(self, file_name, content_preview):
        is_dark = self.page.theme_mode == ft.ThemeMode.DARK
        return ft.Row(
            [
                ft.Container(
                    content=ft.Column(
                        [
                            ft.Text(
                                f"📄 {file_name}",
                                size=15,
                                weight=ft.FontWeight.BOLD,
                                color=ft.colors.WHITE if is_dark else ft.colors.BLUE_800,
                            ),
                            ft.Text(
                                content_preview,
                                selectable=True,
                                size=12,
                                color=ft.colors.WHITE if is_dark else ft.colors.BLACK,
                            )
                        ],
                        spacing=4,
                        tight=True,
                    ),
                    alignment=ft.alignment.center_right,
                    bgcolor=ft.colors.BLUE_800 if is_dark else ft.colors.BLUE_100,
                    padding=ft.padding.symmetric(horizontal=14, vertical=10),
                    border_radius=ft.border_radius.only(
                        top_left=16,
                        top_right=16,
                        bottom_left=16,
                        bottom_right=4,
                    ),
                    margin=ft.margin.only(bottom=6),
                    shadow=ft.BoxShadow(
                        spread_radius=0.5,
                        blur_radius=3,
                        color=ft.colors.with_opacity(0.1, ft.colors.BLACK),
                        offset=ft.Offset(0, 1.5),
                    ),
                )
            ],
            alignment=ft.MainAxisAlignment.END,
        )

    def store_message(self, sender, message):
        """Store a message in the database"""
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        self.cursor.execute('''
            INSERT INTO messages (sender, message, timestamp)
            VALUES (?, ?, ?)
        ''', (sender, message, timestamp))
        self.conn.commit()

    def toggle_theme(self, e):
        """Toggle between light and dark theme"""
        self.page.theme_mode = (
            ft.ThemeMode.DARK
            if self.page.theme_mode == ft.ThemeMode.LIGHT
            else ft.ThemeMode.LIGHT
        )
        self.theme_toggle.icon = (
            ft.icons.LIGHT_MODE
            if self.page.theme_mode == ft.ThemeMode.DARK
            else ft.icons.DARK_MODE
        )
        self.update_theme_colors()
        self.page.update()

    def update_theme_colors(self):
        """Update all color references based on current theme"""
        is_dark = self.page.theme_mode == ft.ThemeMode.DARK
        
        # Update page colors
        self.page.bgcolor = ft.colors.GREY_900 if is_dark else ft.colors.GREY_50
        
        # Update input field
        self.user_input.color = ft.colors.WHITE if is_dark else ft.colors.BLACK
        self.user_input.bgcolor = ft.colors.with_opacity(0.1, ft.colors.GREY_500 if is_dark else ft.colors.GREY_200)
        self.user_input.border_color = ft.colors.GREY_600 if is_dark else ft.colors.GREY_300
        
        # Update buttons
        self.send_button.icon_color = ft.colors.WHITE if is_dark else ft.colors.BLACK
        self.upload_button.icon_color = ft.colors.WHITE if is_dark else ft.colors.BLACK
        self.theme_toggle.icon_color = ft.colors.WHITE if is_dark else ft.colors.BLACK

    def init_ui(self):
        self.page.title = "Bob the lawyer"
        self.page.window_width = 900
        self.page.window_height = 700
        self.page.padding = 20

        # Header with title and theme toggle
        header = ft.Row(
            controls=[
                ft.Text("Chat", size=24, weight=ft.FontWeight.BOLD),
                self.theme_toggle,
            ],
            alignment=ft.MainAxisAlignment.SPACE_BETWEEN,
        )

        main_column = ft.Column(
            expand=True,
            controls=[
                header,
                ft.Divider(color=ft.colors.GREY_600 if self.page.theme_mode == ft.ThemeMode.DARK else ft.colors.GREY_300),
                self.chat,
                ft.Row(
                    [
                        self.upload_button,
                        self.user_input,
                        self.send_button,
                    ],
                    spacing=10,
                    alignment=ft.MainAxisAlignment.START,
                ),
            ],
        )

        content = ft.Row(
            controls=[
                self.sidebar,
                ft.Container(width=1, bgcolor=ft.colors.GREY_600 if self.page.theme_mode == ft.ThemeMode.DARK else ft.colors.GREY_300),
                main_column,
            ],
            expand=True,
        )

        self.page.add(content)
        self.user_input.focus()

    def upload_files(self, e):
        self.file_picker.pick_files(
            allowed_extensions=[
                "pdf", "docx", "pptx", "xls", "xlsx", 
                "jpg", "jpeg", "png", "tiff", "bmp"
            ],
        )
        
    def handle_file_upload(self, e: ft.FilePickerResultEvent):
        if not e.files:
            return
            
        for uploaded_file in e.files:
            file_path = uploaded_file.path
            file_name = uploaded_file.name
            file_ext = os.path.splitext(file_name)[1].lower()
            
            try:
                if file_ext == '.pdf':
                    text = self.extract_text_from_pdf(file_path)
                elif file_ext == '.docx':
                    text = self.extract_text_from_docx(file_path)
                elif file_ext == '.pptx':
                    text = self.extract_text_from_pptx(file_path)
                elif file_ext in ('.xls', '.xlsx'):
                    text = self.extract_text_from_excel(file_path)
                elif file_ext in ('.jpg', '.jpeg', '.png', '.tiff', '.bmp'):
                    text = self.extract_text_from_image(file_path)
                else:
                    text = f"Unsupported file type: {file_ext}"
                
                # Store file content
                self.current_files.append({
                    "name": file_name,
                    "path": file_path,
                    "content": text
                })
                
                # Show file preview as user message
                preview = text[:200] + "..." if len(text) > 200 else text
                self.chat.controls.append(
                    self.create_file_message(file_name, preview))
                
                # Store in database as user message
                self.store_message("user", f"Uploaded document: {file_name}")
                self.store_message("file", f"{file_name}: {text[:200]}...")
                
            except Exception as ex:
                error_msg = f"Error processing {file_name}: {str(ex)}"
                self.chat.controls.append(
                    self.create_bot_message(error_msg))
                    
        self.page.update()

    def extract_text_from_pdf(self, file_path: str) -> str:
        text = []
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text.append(page.extract_text())
        return "\n".join(text)

    def extract_text_from_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    def extract_text_from_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def extract_text_from_excel(self, file_path: str) -> str:
        df = pd.read_excel(file_path)
        return df.to_string()

    def extract_text_from_image(self, file_path: str) -> str:
        try:
            return pytesseract.image_to_string(Image.open(file_path))
        except:
            return f"Image file detected but OCR not available: {file_path}"

    def send_click(self, e):
        question = self.user_input.value.strip()
        if not question and not self.current_files:
            return

        # If there are files, include them in the context
        context = ""
        if self.current_files:
            context = "\n\n[Attached Files Context]\n"
            for file in self.current_files:
                context += f"\nFile: {file['name']}\nContent:\n{file['content'][:1000]}\n"
            self.current_files = []  # Clear files after sending

        if question:
            full_question = question + context
            self.store_message("user", question)
            self.chat.controls.append(self.create_user_message(question))
        else:
            full_question = "Please analyze these documents:" + context
            self.store_message("user", "Uploaded documents for analysis")

        if context:
            self.chat.controls.append(self.create_bot_message("Received documents for analysis"))

        thinking = ft.Container(
            ft.Row([
                ft.ProgressRing(width=20, height=20, stroke_width=2),
                ft.Text("Thinking...")
            ], spacing=10),
            alignment=ft.alignment.center_left,
        )
        self.chat.controls.append(thinking)
        self.user_input.disabled = True
        self.send_button.disabled = True
        self.upload_button.disabled = True
        self.page.update()

        try:
            reply = generate_reply(full_question)
            self.store_message("bot", reply)
        except Exception as err:
            reply = f"⚠️ Error: {str(err)}"
            self.store_message("system", f"Error: {str(err)}")

        self.chat.controls.remove(thinking)
        self.chat.controls.append(self.create_bot_message(reply))

        self.user_input.value = ""
        self.user_input.disabled = False
        self.send_button.disabled = False
        self.upload_button.disabled = False
        self.page.update()
        self.user_input.focus()

    def __del__(self):
        """Close database connection when the app is closed"""
        if hasattr(self, 'conn'):
            self.conn.close()


def main(page: ft.Page):
    LawyerChatBotApp(page)

ft.app(target=main)